from pptx import Presentation
from pptx.util import Inches
from knowledge.schema import PaperKnowledge
from pathlib import Path
from datetime import datetime

def generate_ppt(paper: PaperKnowledge, summary: str, flaws: str, comparison: str, user_id: int = None) -> str:
    """Generate professional PPT and save it correctly"""
    
    # Create user-specific folder if user_id is provided
    base_dir = Path("user_data")
    if user_id:
        ppt_dir = base_dir / str(user_id) / "ppts"
    else:
        ppt_dir = base_dir / "ppts"
    
    ppt_dir.mkdir(parents=True, exist_ok=True)

    # Proper filename as string
    safe_title = paper.title[:40].replace(" ", "_").replace("/", "_")
    filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    filepath = ppt_dir / filename

    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = paper.title
    if paper.authors:
        slide.placeholders[1].text = ", ".join(paper.authors[:5])

    # Summary Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = summary[:800] if summary else "Summary not available"

    # Insights / Key Findings
    if flaws:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Flaws & Limitations"
        tf = slide.placeholders[1].text_frame
        tf.text = flaws[:700]

    # Comparison
    if comparison:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Comparison with Related Work"
        tf = slide.placeholders[1].text_frame
        tf.text = comparison[:700]

    prs.save(str(filepath))   # Ensure it's a string path
    print(f"✅ PPT saved: {filepath}")
    return str(filepath)